"""Leitura de documentos em múltiplos formatos: PDF, Word, Excel, PowerPoint, Markdown, CSV, JSON e HTML."""
import csv
import json
import os

SUPPORTED_EXTS = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".xlsx": "Excel",
    ".pptx": "PowerPoint",
    ".md": "Markdown",
    ".csv": "CSV",
    ".json": "JSON",
    ".html": "HTML",
}


def _ler_docx(caminho: str) -> str:
    from docx import Document

    doc = Document(caminho)
    linhas = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            linha_tabela = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if linha_tabela:
                linhas.append(linha_tabela)
    return "\n".join(linhas)


def _ler_pptx(caminho: str) -> str:
    from pptx import Presentation

    prs = Presentation(caminho)
    blocos = []
    for i, slide in enumerate(prs.slides, 1):
        blocos.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for par in shape.text_frame.paragraphs:
                    texto = "".join(run.text for run in par.runs)
                    if texto.strip():
                        blocos.append(texto)
            elif shape.has_table:
                for row in shape.table.rows:
                    linha = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if linha:
                        blocos.append(linha)
    return "\n".join(blocos)


def _ler_xlsx(caminho: str) -> str:
    import pandas as pd

    xls = pd.ExcelFile(caminho)
    partes = []
    for sheet in xls.sheet_names:
        df = xls.parse(sheet, dtype=str).fillna("")
        partes.append(f"--- Planilha: {sheet} ---\n" + df.to_csv(index=False))
    return "\n".join(partes)


def _ler_csv(caminho: str) -> str:
    with open(caminho, newline="", encoding="utf-8") as f:
        return f.read()


def _ler_json(caminho: str) -> str:
    with open(caminho, encoding="utf-8") as f:
        return json.dumps(json.load(f), ensure_ascii=False, indent=2)


def _ler_html(caminho: str) -> str:
    from html.parser import HTMLParser

    class TextoHTML(HTMLParser):
        def __init__(self):
            super().__init__()
            self.pedacos = []

        def handle_data(self, data):
            if data.strip():
                self.pedacos.append(data.strip())

    parser = TextoHTML()
    with open(caminho, encoding="utf-8") as f:
        parser.feed(f.read())
    return " ".join(parser.pedacos)


def _ler_markdown(caminho: str) -> str:
    with open(caminho, encoding="utf-8") as f:
        return f.read()


def load_document(caminho: str) -> str:
    """Extrai o texto de um documento. Levanta ValueError para extensão não suportada."""
    ext = os.path.splitext(caminho)[1].lower()
    if ext not in SUPPORTED_EXTS:
        raise ValueError(f"Formato não suportado: {ext}. Suportados: {', '.join(sorted(SUPPORTED_EXTS))}")

    if ext == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(caminho)
        return "\n".join(pagina.extract_text() or "" for pagina in reader.pages)
    if ext == ".docx":
        return _ler_docx(caminho)
    if ext == ".pptx":
        return _ler_pptx(caminho)
    if ext == ".xlsx":
        return _ler_xlsx(caminho)
    if ext == ".csv":
        return _ler_csv(caminho)
    if ext == ".json":
        return _ler_json(caminho)
    if ext == ".html":
        return _ler_html(caminho)
    if ext == ".md":
        return _ler_markdown(caminho)
    raise ValueError(f"Formato não suportado: {ext}")