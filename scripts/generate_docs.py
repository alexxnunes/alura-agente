"""Gera a base documental corporativa completa para os 6 domínios/segmentos do Challenge Alura Agentes:
1. Loja Online / E-commerce
2. SaaS / Plataforma Digital
3. Empresa de Logística / Envios
4. Clínica de Saúde / Consultório Médico
5. Plataforma Educativa / Escola Online
6. Fintech / Banco Digital

Cobre todos os 8 formatos de arquivo suportados: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), Markdown (.md), CSV (.csv), JSON (.json) e HTML (.html).
"""
import csv
import json
import os
import pandas as pd
from docx import Document as DocxDocument
from fpdf import FPDF, XPos, YPos
from pptx import Presentation
from pptx.util import Inches

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS_DIR = os.path.join(BASE_DIR, "data", "docs")


# ==========================================
# 1. LOJA ONLINE / E-COMMERCE
# ==========================================
def gerar_ecommerce(dir_path: str):
    os.makedirs(dir_path, exist_ok=True)

    # 1.1 Politica de privacidade (.md)
    md_content = """# Política de Privacidade — VendaMax Brasil E-commerce

Última atualização: 15 de janeiro de 2026.

## 1. Informações Coletadas
A VendaMax coleta os seguintes dados fornecidos diretamente pelo cliente:
- **Dados Cadastrais**: Nome completo, CPF, endereço de entrega e faturamento, e-mail e telefone de contato.
- **Dados de Navegação**: Endereço IP, cookies de sessão, histórico de visualização de produtos e preferências de compra.
- **Dados de Pagamento**: Informações criptografadas processadas diretamente por intermediadores homologados (Stripe e Wirecard). A VendaMax não armazena números completos de cartão de crédito.

## 2. Finalidade do Tratamento
Os dados coletados destinam-se exclusivamente a:
- Processar e faturar pedidos de compra e emissão da Nota Fiscal Eletrônica (NF-e).
- Entregar produtos no endereço indicado e realizar atendimento pós-venda.
- Prevenir fraudes e garantir a segurança das transações comerciais.
- Enviar comunicações promocionais e status do pedido, com opção de descadastramento (opt-out) a qualquer momento.

## 3. Compartilhamento de Dados
Os dados são compartilhados apenas com parceiros essenciais para a operação:
- Transportadoras e operadoras logísticas credenciadas para cumprimento da entrega.
- Gateways de pagamento seguros para autorização bancária.
- Órgãos públicos e autoridades fiscais conforme exigência da legislação tributária brasileira.

## 4. Direitos do Titular e Contato do DPO
Em conformidade com a LGPD (Lei 13.709/2018), o cliente pode solicitar a confirmação, correção, portabilidade ou eliminação de seus dados pelo e-mail **dpo@vendamax.com.br**. As solicitações são respondidas em até 15 dias corridos.
"""
    with open(os.path.join(dir_path, "politica_privacidade.md"), "w", encoding="utf-8") as f:
        f.write(md_content)

    # 1.2 Politica de reembolso e devolucoes (.pdf)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", "B", 15)
    pdf.cell(0, 10, "Politica de Reembolso e Devolucoes - VendaMax", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(3)
    pdf.set_font("helvetica", "", 10)
    regras_ecommerce = [
        "1. Direito de Arrependimento (Art. 49 do CDC):",
        "   - O cliente tem ate 7 (sete) dias corridos apos o recebimento para solicitar o cancelamento.",
        "   - O frete de devolucao e 100% gratuito atraves de codigo de postagem reversa dos Correios.",
        "",
        "2. Troca por Defeito ou Avaria:",
        "   - Produtos nao duraveis: prazo de 30 dias para reclamacao.",
        "   - Produtos duraveis (eletronicos, eletrodomesticos): prazo de 90 dias conforme CDC.",
        "   - Laudo tecnico realizado em ate 5 dias uteis apos o recebimento do item em nosso CD.",
        "",
        "3. Formas e Prazos de Restituicao de Valores:",
        "   - Pix: estorno em ate 24 horas uteis na mesma conta bancaria de origem.",
        "   - Cartao de Credito: notificacao da operadora em ate 3 dias uteis; estorno na fatura atual ou seguinte.",
        "   - Boleto Bancario: transferencia TED/Pix em ate 3 dias uteis apos envio dos dados bancarios.",
        "",
        "4. Condicoes do Produto para Aceite:",
        "   - O item deve estar acompanhado da 1a via da Nota Fiscal (DANFE).",
        "   - Produto na embalagem original, com todos os lacres, manuais e acessorios sem avarias.",
    ]
    for r in regras_ecommerce:
        pdf.multi_cell(0, 5, r, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.output(os.path.join(dir_path, "politica_reembolso_devolucoes.pdf"))

    # 1.3 FAQ (.json)
    faq_data = {
        "dominio": "Loja Online / E-commerce",
        "empresa": "VendaMax Brasil",
        "perguntas_frequentes": [
            {
                "id": "FAQ-EC-01",
                "pergunta": "Quais são as formas de pagamento aceitas?",
                "resposta": "Aceitamos Pix com 5% de desconto à vista, cartão de crédito em até 10x sem juros (bandeiras Visa, Mastercard, Elo, Hipercard e American Express) e boleto bancário à vista."
            },
            {
                "id": "FAQ-EC-02",
                "pergunta": "Existe cupom de desconto para primeira compra?",
                "resposta": "Sim! Novos clientes podem utilizar o cupom 'BEMVINDO10' no carrinho para obter 10% de desconto em pedidos acima de R$ 100,00."
            },
            {
                "id": "FAQ-EC-03",
                "pergunta": "Como posso alterar o endereço de entrega do meu pedido?",
                "resposta": "A alteração de endereço só pode ser feita se a Nota Fiscal ainda não tiver sido emitida. Caso a NF já tenha sido gerada, por normas fiscais não é possível alterar o destino, sendo necessário cancelar e refazer a compra."
            },
            {
                "id": "FAQ-EC-04",
                "pergunta": "O que fazer se meu pedido atrasar?",
                "resposta": "Se o prazo de entrega expirou, entre em contato imediatamente com nosso SAC pelo WhatsApp (11) 98765-4321 ou e-mail sac@vendamax.com.br informando o número do pedido. Abriremos um protocolo de urgência com a transportadora."
            }
        ]
    }
    with open(os.path.join(dir_path, "faq.json"), "w", encoding="utf-8") as f:
        json.dump(faq_data, f, ensure_ascii=False, indent=2)

    # 1.4 Guia de envios e entregas (.docx)
    doc = DocxDocument()
    doc.add_heading("Guia de Envios e Prazos de Entrega — VendaMax", level=1)
    doc.add_paragraph("Este guia detalha os procedimentos logísticos, custos e prazos de entrega praticados pela VendaMax.")
    
    doc.add_heading("1. Modalidades de Envio", level=2)
    p1 = doc.add_paragraph()
    p1.add_run("• Sedex / Expresso: ").bold = True
    p1.add_run("Prazo de 1 a 3 dias úteis para capitais e regiões metropolitanas.")
    p2 = doc.add_paragraph()
    p2.add_run("• PAC / Encomenda Padrão: ").bold = True
    p2.add_run("Prazo de 4 a 8 dias úteis para o interior e regiões periféricas.")
    p3 = doc.add_paragraph()
    p3.add_run("• Entrega Super Expressa (Same Day): ").bold = True
    p3.add_run("Pedidos confirmados até as 12h são entregues até as 21h do mesmo dia na Grande São Paulo e Rio de Janeiro.")

    doc.add_heading("2. Política de Frete Grátis", level=2)
    doc.add_paragraph("• Sul e Sudeste: Frete grátis em compras acima de R$ 199,00.")
    doc.add_paragraph("• Centro-Oeste e Nordeste: Frete grátis em compras acima de R$ 299,00.")
    doc.add_paragraph("• Região Norte: Frete grátis em compras selecionadas acima de R$ 399,00.")

    doc.add_heading("3. Tentativas de Entrega", level=2)
    doc.add_paragraph("A transportadora realiza até 3 (três) tentativas de entrega em dias úteis consecutivos. Caso não haja ninguém no local, a encomenda permanece disponível na agência ou ponto de coleta mais próximo por até 7 dias corridos antes de retornar ao centro de distribuição.")
    doc.save(os.path.join(dir_path, "guia_envios_entregas.docx"))

    # 1.5 Termos e condicoes (.html)
    html_content = """<!DOCTYPE html>
<html lang="pt-BR">
<head>
    <meta charset="UTF-8">
    <title>Termos e Condições de Uso - VendaMax</title>
</head>
<body>
    <h1>Termos e Condições Gerais de Uso e Compra</h1>
    <p>Bem-vindo à <strong>VendaMax Brasil Comércio Eletrônico Ltda</strong> (CNPJ 12.345.678/0001-90).</p>
    
    <h2>1. Objeto</h2>
    <p>Estes termos regem o uso do site e a compra de produtos comercializados pela plataforma VendaMax por usuários consumidores.</p>
    
    <h2>2. Cadastro e Responsabilidade do Usuário</h2>
    <p>O usuário declara ter idade legal mínima de 18 anos ou ser emancipado. O usuário é o único responsável pela guarda e confidencialidade de sua senha de acesso.</p>
    
    <h2>3. Preços e Disponibilidade de Estoque</h2>
    <p>Os preços promocionais, cupons e condições de pagamento são válidos apenas durante o período divulgado ou enquanto durarem os estoques físicos. Em caso de divergência de preço no carrinho, prevalecerá o valor final discriminado no checkout.</p>
    
    <h2>4. Propriedade Intelectual e Foro</h2>
    <p>Todo o conteúdo textual, fotográfico e marcas do site pertencem à VendaMax. Fica eleito o Foro da Comarca de São Paulo/SP para dirimir quaisquer controvérsias decorrentes destes Termos.</p>
</body>
</html>"""
    with open(os.path.join(dir_path, "termos_condicoes.html"), "w", encoding="utf-8") as f:
        f.write(html_content)


# ==========================================
# 2. SAAS / PLATAFORMA DIGITAL
# ==========================================
def gerar_saas(dir_path: str):
    os.makedirs(dir_path, exist_ok=True)

    # 2.1 Base de conhecimento (.md)
    md_content = """# Base de Conhecimento — CloudSync Pro Platform

Bem-vindo à documentação técnica oficial da plataforma **CloudSync Pro**, a solução corporativa para automação de dados e integrações em nuvem.

## 1. Arquitetura e Módulos Principais
A plataforma CloudSync Pro é composta por quatro módulos integrados:
1. **Pipeline Builder**: Criação visual de fluxos de transformação ETL/ELT sem necessidade de código (no-code / low-code).
2. **API Management & Webhooks**: Gateway para recepção e disparo de eventos HTTP em tempo real com retry exponencial automático em caso de falha 5xx.
3. **Audit Log & Compliance**: Rastreabilidade completa de todas as ações de usuários, alterações de configuração e logs de execução armazenados por até 365 dias.
4. **Connector Hub**: Mais de 120 conectores nativos pré-construídos para bancos de dados (PostgreSQL, MySQL, Oracle, MongoDB) e ferramentas SaaS (Salesforce, HubSpot, Jira, Slack).

## 2. Autenticação e Segurança
- **API Keys**: Geradas com permissões granulares por workspace (Read-only, Write, Admin).
- **OAuth 2.0 / SAML 2.0**: Suporte a Single Sign-On (SSO) com provedores Okta, Azure Active Directory e Google Workspace no plano Enterprise.
- **Rate Limits**: 1.000 requisições por minuto no plano Pro e 5.000 requisições por minuto no plano Enterprise.
"""
    with open(os.path.join(dir_path, "base_conhecimento.md"), "w", encoding="utf-8") as f:
        f.write(md_content)

    # 2.2 FAQ Suporte (.json)
    faq_data = {
        "dominio": "SaaS / Plataforma Digital",
        "produto": "CloudSync Pro",
        "faq_suporte": [
            {
                "id": "FAQ-SAAS-01",
                "pergunta": "Como redefinir a autenticação em dois fatores (2FA)?",
                "resposta": "Para redefinir o 2FA, acesse 'Configurações de Perfil' > 'Segurança' e clique em 'Recuperar Chave de Segurança'. Se tiver perdido o acesso ao autenticador, solicite a recuperação ao Administrador do Workspace."
            },
            {
                "id": "FAQ-SAAS-02",
                "pergunta": "Quais são os níveis de SLA de atendimento por plano?",
                "resposta": "Plano Starter: até 24 horas úteis por e-mail; Plano Pro: até 4 horas úteis via chat e e-mail; Plano Enterprise: até 1 hora de resposta 24/7 com gerente técnico de contas dedicado."
            },
            {
                "id": "FAQ-SAAS-03",
                "pergunta": "Como exportar relatórios de execução de pipelines?",
                "resposta": "No menu 'Analytics', selecione o período desejado e clique no botão 'Exportar' no canto superior direito. Os dados podem ser baixados nos formatos CSV, JSON ou PDF."
            },
            {
                "id": "FAQ-SAAS-04",
                "pergunta": "Qual é a garantia de disponibilidade (Uptime) do sistema?",
                "resposta": "Oferecemos garantia de 99.9% de uptime para clientes Pro e 99.95% para clientes Enterprise, monitorado publicamente na página de status (status.cloudsyncpro.io)."
            }
        ]
    }
    with open(os.path.join(dir_path, "faq_suporte.json"), "w", encoding="utf-8") as f:
        json.dump(faq_data, f, ensure_ascii=False, indent=2)

    # 2.3 Politica de privacidade (.docx)
    doc = DocxDocument()
    doc.add_heading("Política de Privacidade e Proteção de Dados — CloudSync Pro", level=1)
    doc.add_paragraph("A CloudSync Technologies Inc. adota rigorosos padrões internacionais de segurança da informação (ISO/IEC 27001 e SOC 2 Tipo II).")
    
    doc.add_heading("1. Dados Processados pela Plataforma", level=2)
    doc.add_paragraph("Processamos apenas dados estritamente necessários para o funcionamento das integrações configuradas pelo cliente. Todos os dados em trânsito são criptografados com TLS 1.3 e em repouso com AES-256 bits.")
    
    doc.add_heading("2. Isolamento de Ambientes (Multi-tenancy)", level=2)
    doc.add_paragraph("Cada workspace possui esquema de dados isolado com chaves de criptografia segregadas. Nenhum dado do cliente é utilizado para treinamento de modelos de IA sem consentimento expresso por contrato.")
    
    doc.add_heading("3. Retenção e Backups", level=2)
    doc.add_paragraph("Backups automáticos são realizados a cada 6 horas com retenção de 30 dias em regiões geograficamente redundantes. Mediante encerramento do contrato, todos os dados do cliente são destruídos de forma segura em até 30 dias úteis.")
    doc.save(os.path.join(dir_path, "politica_privacidade.docx"))

    # 2.4 Planos e precos (.xlsx)
    df_planos = pd.DataFrame([
        {"Plano": "Starter", "Preco_Mensal_BRL": 49.00, "Preco_Anual_BRL": 470.00, "Usuarios_Inclusos": 3, "Pipelines_Ativos": 10, "Armazenamento_GB": 10, "SLA_Suporte": "24h úteis", "SSO_SAML": "Não"},
        {"Plano": "Pro", "Preco_Mensal_BRL": 149.00, "Preco_Anual_BRL": 1430.00, "Usuarios_Inclusos": 15, "Pipelines_Ativos": 50, "Armazenamento_GB": 100, "SLA_Suporte": "4h úteis", "SSO_SAML": "Opcional"},
        {"Plano": "Enterprise", "Preco_Mensal_BRL": 499.00, "Preco_Anual_BRL": 4790.00, "Usuarios_Inclusos": 9999, "Pipelines_Ativos": 9999, "Armazenamento_GB": 1000, "SLA_Suporte": "1h 24/7", "SSO_SAML": "Incluso"},
    ])
    with pd.ExcelWriter(os.path.join(dir_path, "planos_e_precos.xlsx"), engine="openpyxl") as writer:
        df_planos.to_excel(writer, sheet_name="Tabela_Planos", index=False)

    # 2.5 Termos de uso (.pdf)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", "B", 15)
    pdf.cell(0, 10, "Termos de Uso do Software SaaS - CloudSync Pro", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(3)
    pdf.set_font("helvetica", "", 10)
    termos_saas = [
        "1. Concessao de Licenca de Uso:",
        "   A CloudSync concede licenca nao exclusiva, revogavel e intransferivel para uso da plataforma.",
        "",
        "2. Nivel de Servico (SLA) e Compensacao Financeira:",
        "   - Disponibilidade minima mensal garantida: 99.9% para Pro e 99.95% para Enterprise.",
        "   - Caso o SLA nao seja atingido, o cliente recebera creditos proporcionais de 10% a 30% na fatura.",
        "",
        "3. Cancelamento e Rescisao:",
        "   - Planos mensais podem ser cancelados a qualquer momento sem multa com efeito no ciclo seguinte.",
        "   - Planos anuais com desconto preveem aviso previo de 30 dias com multa rescisoria de 20%.",
        "",
        "4. Restricoes de Uso:",
        "   E vedado praticar engenharia reversa, sublicenciar ou utilizar a API para envio de spams ou ataques.",
    ]
    for t in termos_saas:
        pdf.multi_cell(0, 5, t, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.output(os.path.join(dir_path, "termos_de_uso.pdf"))


# ==========================================
# 3. EMPRESA DE LOGÍSTICA / ENVIOS
# ==========================================
def gerar_logistica(dir_path: str):
    os.makedirs(dir_path, exist_ok=True)

    # 3.1 Politica de envios (.pdf)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", "B", 15)
    pdf.cell(0, 10, "Politica de Envios e Cargas - TransLogistica Express", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(3)
    pdf.set_font("helvetica", "", 10)
    regras_log = [
        "1. Limites de Dimensoes e Peso por Volume:",
        "   - Peso maximo por encomenda individual: 30 kg.",
        "   - Dimensao maxima de comprimento: 105 cm.",
        "   - Soma das dimensoes (Comprimento + Largura + Altura): maximo de 200 cm.",
        "",
        "2. Itens Proibidos e Nao Transportaveis:",
        "   - Explosivos, gases comprimidos, liquidos inflamaveis e substancias toxicas.",
        "   - Animais vivos, restos mortais, joias preciosas e dinheiro em especie.",
        "   - Cargas pereciveis sem contratacao especifica de transporte com cadeia refrigerada.",
        "",
        "3. Horarios de Corte (Cut-off) para Embarque:",
        "   - Solicitacao de coleta no mesmo dia: ate as 14h00.",
        "   - Despacho em centro de distribuicao: ate as 17h30 para embarque noturno.",
    ]
    for r in regras_log:
        pdf.multi_cell(0, 5, r, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.output(os.path.join(dir_path, "politica_envios.pdf"))

    # 3.2 Procedimento de rastreamento (.md)
    md_content = """# Procedimento de Rastreamento de Pedidos — TransLogística Express

## 1. Estrutura do Código de Rastreamento
Todas as encomendas expedidas pela TransLogística recebem um código alfanumérico padrão de 11 caracteres:
- Formato: `TL-` + 7 dígitos numéricos + `BR` (Exemplo: `TL-8392014BR`).

## 2. Etapas e Status de Rastreio
O sistema é atualizado automaticamente a cada ponto de conferência por leitor ótico:
1. **Objeto Coletado / Postado**: Carga recebida na filial de origem e bipada no sistema.
2. **Em Transferência (Transfer Hub)**: Carga em trânsito rodoviário ou aéreo entre centros de triagem.
3. **Recebido no Centro de Distribuição de Destino**: Carga em triagem para rota urbana.
4. **Saiu para Entrega ao Destinatário**: Pacote embarcado no veículo de entrega porta a porta.
5. **Entregue ao Destinatário**: Finalizado com captura da assinatura digital e foto do comprovante.

## 3. Frequência de Atualização
As atualizações de status ocorrem em tempo real e são disponibilizadas no portal e aplicativo em no máximo **30 minutos** após a leitura no ponto de controle.
"""
    with open(os.path.join(dir_path, "rastreamento_pedidos.md"), "w", encoding="utf-8") as f:
        f.write(md_content)

    # 3.3 Reembolsos e sinistros (.docx)
    doc = DocxDocument()
    doc.add_heading("Política de Reembolsos e Sinistros de Carga — TransLogística", level=1)
    doc.add_paragraph("Normas para abertura de ocorrências, coberturas securitárias e prazos de indenização.")
    
    doc.add_heading("1. Cobertura do Seguro de Transporte", level=2)
    doc.add_paragraph("Todas as encomendas possuem seguro obrigatório RCTR-C contra acidentes, tombamento e roubo com cobertura automática de até R$ 10.000,00 por declaração de valor na Nota Fiscal.")
    
    doc.add_heading("2. Prazos para Abertura de Sinistro", level=2)
    doc.add_paragraph("• Extravio ou Roubo de Carga: abertura automática após 48 horas sem atualização de status ou mediante boletim de ocorrência.")
    doc.add_paragraph("• Avaria / Quebra de Mercadoria: o destinatário ou remetente deve formalizar a reclamação em até 15 dias corridos após a data de entrega, com envio de fotos da embalagem e do item avariado.")
    
    doc.add_heading("3. Prazo de Pagamento da Indenização", level=2)
    doc.add_paragraph("Após a conclusão do laudo de vistoria (realizado em até 5 dias úteis), o pagamento da indenização do valor total constante na NF é realizado via transferência bancária em até 10 (dez) dias úteis.")
    doc.save(os.path.join(dir_path, "reembolsos_sinistros.docx"))

    # 3.4 FAQ Logistica (.json)
    faq_data = {
        "dominio": "Empresa de Logística / Envios",
        "empresa": "TransLogística Express",
        "perguntas_frequentes": [
            {
                "id": "FAQ-LOG-01",
                "pergunta": "Quantas tentativas de entrega são realizadas?",
                "resposta": "Realizamos 3 (três) tentativas de entrega em dias consecutivos. Se não houver sucesso, o pacote aguarda retirada na filial mais próxima por 5 dias antes de iniciar o processo de devolução ao remetente."
            },
            {
                "id": "FAQ-LOG-02",
                "pergunta": "Posso alterar o endereço após a carga sair para entrega?",
                "resposta": "Não é possível alterar o endereço de um pacote já em rota de entrega no mesmo dia. A solicitação de redirecionamento deve ser feita pelo remetente antes da saída do centro de distribuição."
            },
            {
                "id": "FAQ-LOG-03",
                "pergunta": "Como comprovar que o pacote foi entregue?",
                "resposta": "Nosso motorista colhe a assinatura digital, documento (RG/CPF) do recebedor e foto da fachada/comprovante geolocalizado, disponíveis para consulta no link do rastreador."
            }
        ]
    }
    with open(os.path.join(dir_path, "faq_logistica.json"), "w", encoding="utf-8") as f:
        json.dump(faq_data, f, ensure_ascii=False, indent=2)

    # 3.5 Reclamacoes e SAC (.html)
    html_content = """<!DOCTYPE html>
<html lang="pt-BR">
<head>
    <meta charset="UTF-8">
    <title>Canais de Atendimento e Reclamações - TransLogística</title>
</head>
<body>
    <h1>Processo de Reclamações e Atendimento ao Cliente</h1>
    <p>Central de Atendimento ao Cliente da <strong>TransLogística Express</strong>.</p>
    
    <h2>1. Canais Oficiais de Atendimento</h2>
    <ul>
        <li><strong>Telefone SAC 0800</strong>: 0800-777-8900 (Segunda a Sexta, das 08h00 às 20h00 e Sábados das 08h00 às 14h00).</li>
        <li><strong>WhatsApp Corporativo</strong>: (11) 97777-8900 (Atendimento automatizado 24h e atendentes humanos em horário comercial).</li>
        <li><strong>Portal de Chamados</strong>: suporte.translogistica.com.br.</li>
    </ul>
    
    <h2>2. SLAs de Resolução</h2>
    <p>Todo chamado aberto recebe um protocolo único de 8 dígitos. O prazo de primeira resposta técnica é de até <strong>2 horas úteis</strong>, e a resolução definitiva de ocorrências ocorre em no máximo <strong>48 horas úteis</strong>.</p>
</body>
</html>"""
    with open(os.path.join(dir_path, "reclamacoes_atendimento.html"), "w", encoding="utf-8") as f:
        f.write(html_content)


# ==========================================
# 4. CLÍNICA DE SAÚDE / CONSULTÓRIO MÉDICO
# ==========================================
def gerar_saude(dir_path: str):
    os.makedirs(dir_path, exist_ok=True)

    # 4.1 Privacidade de dados do paciente (.pdf)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", "B", 15)
    pdf.cell(0, 10, "Politica de Privacidade de Dados de Pacientes - Vida & Saude", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(3)
    pdf.set_font("helvetica", "", 10)
    regras_saude = [
        "1. Conformidade com a LGPD e Resolucao CFM no 1.821/2007:",
        "   - Dados de saude sao classificados como dados pessoais sensiveis (Art. 11 da LGPD).",
        "   - O prontuario medico eletronico e protegido por criptografia de nivel hospitalar.",
        "",
        "2. Guarda e Retencao de Prontuarios:",
        "   - Conforme legislacao federal, os prontuarios medicos sao arquivados pelo prazo minimo de 20 anos.",
        "   - O paciente tem o direito de solicitar copia integral do seu prontuario a qualquer momento.",
        "",
        "3. Sigilo Medico e Acesso Restrito:",
        "   - O acesso aos dados clinicos e restrito exclusivamente a equipe medica e assistencial responsavel.",
        "   - E estritamente proibido o compartilhamento de informacoes clinicas para fins comerciais ou publicitarios.",
    ]
    for r in regras_saude:
        pdf.multi_cell(0, 5, r, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.output(os.path.join(dir_path, "privacidade_dados_paciente.pdf"))

    # 4.2 FAQ Consultas (.json)
    faq_data = {
        "dominio": "Clínica de Saúde / Consultório Médico",
        "clinica": "Clínica Vida & Saúde Integrada",
        "faq_consultas_agendamentos": [
            {
                "id": "FAQ-SAUDE-01",
                "pergunta": "Qual é o horário de funcionamento da clínica?",
                "resposta": "A clínica funciona de segunda a sexta-feira, das 07h00 às 21h00, e aos sábados, das 07h30 às 14h00."
            },
            {
                "id": "FAQ-SAUDE-02",
                "pergunta": "Como funciona o agendamento de telemedicina?",
                "resposta": "As teleconsultas podem ser agendadas pelo nosso portal ou WhatsApp. O link seguro da videochamada é enviado por e-mail e SMS 30 minutos antes do atendimento. Receitas e atestados possuem assinatura digital com certificado ICP-Brasil."
            },
            {
                "id": "FAQ-SAUDE-03",
                "pergunta": "Qual é o prazo para retorno de consulta sem custo adicional?",
                "resposta": "O retorno médico para apresentação de exames solicitados na consulta inicial tem validade de até 21 dias corridos a contar da data do atendimento original."
            }
        ]
    }
    with open(os.path.join(dir_path, "faq_consultas_agendamentos.json"), "w", encoding="utf-8") as f:
        json.dump(faq_data, f, ensure_ascii=False, indent=2)

    # 4.3 Cancelamentos e reagendamentos (.md)
    md_content = """# Política de Cancelamentos e Reagendamentos — Clínica Vida & Saúde

## 1. Antecedência Mínima
- O paciente pode cancelar ou reagendar consultas e exames sem qualquer custo ou penalidade com antecedência mínima de **24 horas úteis**.
- O reagendamento pode ser realizado pelo aplicativo da clínica, pelo telefone (11) 3456-7890 ou WhatsApp oficial.

## 2. Política de No-Show (Ausência sem Aviso Prévio)
- Em consultas particulares, ausências sem comunicação prévia de 24 horas implicam a retenção de **30% do valor da consulta** a título de taxa de reserva de agenda médica para remarcação.
- Em consultas por convênio, o não comparecimento reiterado (mais de 2 ausências sem aviso) bloqueia agendamentos futuros online, exigindo confirmação presencial.

## 3. Tolerância de Horário
A tolerância máxima para atraso do paciente é de **15 minutos**. Após esse período, o atendimento fica condicionado à disponibilidade de encaixe médico para não prejudicar os demais pacientes agendados.
"""
    with open(os.path.join(dir_path, "cancelamentos_reagendamentos.md"), "w", encoding="utf-8") as f:
        f.write(md_content)

    # 4.4 Convenios e coberturas (.xlsx)
    df_convenios = pd.DataFrame([
        {"Convenio": "Unimed Nacional", "Planos_Aceitos": "Básico, Especial, Master", "Consultas": "Sim", "Exames_Laboratoriais": "Sim", "Ultrassom_ECG": "Sim", "Exige_Guia_Previa": "Não"},
        {"Convenio": "Bradesco Saúde", "Planos_Aceitos": "Top Nacional, Premium", "Consultas": "Sim", "Exames_Laboratoriais": "Sim", "Ultrassom_ECG": "Sim", "Exige_Guia_Previa": "Não"},
        {"Convenio": "Amil", "Planos_Aceitos": "Amil Fácil, Amil 400 a 700", "Consultas": "Sim", "Exames_Laboratoriais": "Sim", "Ultrassom_ECG": "Sim", "Exige_Guia_Previa": "Sim (alguns exames)"},
        {"Convenio": "SulAmérica", "Planos_Aceitos": "Exato, Especial 100, Executivo", "Consultas": "Sim", "Exames_Laboratoriais": "Sim", "Ultrassom_ECG": "Sim", "Exige_Guia_Previa": "Não"},
        {"Convenio": "NotreDame Intermédica", "Planos_Aceitos": "Smart, Advance, Premium", "Consultas": "Sim", "Exames_Laboratoriais": "Sim", "Ultrassom_ECG": "Sim", "Exige_Guia_Previa": "Sim"},
    ])
    with pd.ExcelWriter(os.path.join(dir_path, "convenios_coberturas.xlsx"), engine="openpyxl") as writer:
        df_convenios.to_excel(writer, sheet_name="Convenios_Coberturas", index=False)

    # 4.5 Instrucoes pre e pos consulta (.docx)
    doc = DocxDocument()
    doc.add_heading("Instruções Pré e Pós-Consulta e Exames — Vida & Saúde", level=1)
    
    doc.add_heading("1. Preparo para Exames Laboratoriais", level=2)
    doc.add_paragraph("• Exames de sangue (Perfil lipídico, colesterol e triglicérides): Jejum absoluto de 12 horas (água é permitida com moderação).")
    doc.add_paragraph("• Glicemia de jejum: Jejum de 8 a 10 horas.")
    doc.add_paragraph("• Não ingerir bebidas alcoólicas nas 72 horas que antecedem a coleta.")

    doc.add_heading("2. Documentos Obrigatórios para a Consulta", level=2)
    doc.add_paragraph("• Documento oficial com foto (RG, CNH ou Passaporte) e cartão do convênio dentro da validade.")
    doc.add_paragraph("• Exames laboratoriais e laudos de imagem realizados nos últimos 6 meses.")
    doc.add_paragraph("• Lista com nomes e dosagens de todos os medicamentos de uso contínuo.")
    doc.save(os.path.join(dir_path, "instrucoes_pre_pos_consulta.docx"))


# ==========================================
# 5. PLATAFORMA EDUCATIVA / ESCOLA ONLINE
# ==========================================
def gerar_educacao(dir_path: str):
    os.makedirs(dir_path, exist_ok=True)

    # 5.1 Regulamento do estudante (.pdf)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", "B", 15)
    pdf.cell(0, 10, "Regulamento Academico do Estudante - Alura Tech Academy", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(3)
    pdf.set_font("helvetica", "", 10)
    regras_edu = [
        "1. Direitos do Estudante:",
        "   - Acesso irrestrito 24/7 a todas as videoaulas, exercicios e projetos praticos da trilha contratada.",
        "   - Participacao ativa no forum de duvidas com suporte de instrutores e monitores qualificados.",
        "",
        "2. Criterios de Aprovacao e Emissao de Certificados:",
        "   - Conclusao obrigatoria de 100% dos modulos em video da formacao.",
        "   - Nota minima de aproveitamento de 75% nas avaliacoes de multipla escolha e projetos praticos.",
        "   - Certificados emitidos digitalmente com codigo hash e QR Code de autenticidade reconhecido.",
        "",
        "3. Codigo de Etica e Integridade Academica:",
        "   - E expressamente proibido o compartilhamento de credenciais de acesso com terceiros.",
        "   - Casos comprovados de plagio em projetos de conclusao resultam em reprovacao sem direito a certificado.",
    ]
    for r in regras_edu:
        pdf.multi_cell(0, 5, r, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.output(os.path.join(dir_path, "regulamento_estudante.pdf"))

    # 5.2 Politica de reembolso (.docx)
    doc = DocxDocument()
    doc.add_heading("Política de Reembolso e Cancelamento de Matrícula — Alura Tech", level=1)
    
    doc.add_heading("1. Garantia Incondicional de 7 Dias", level=2)
    doc.add_paragraph("O aluno tem direito a cancelar sua matrícula e receber 100% de estorno do valor pago em até 7 (sete) dias corridos a partir da data de confirmação da compra, sem necessidade de justificativa.")
    
    doc.add_heading("2. Cancelamento após o Período de Garantia", level=2)
    doc.add_paragraph("• Assinaturas Mensais: o cancelamento interrompe a cobrança do mês seguinte, mantendo o acesso até o fim do ciclo mensal vigente.")
    doc.add_paragraph("• Planos Anuais: o cancelamento antecipado após os 7 dias acarreta taxa administrativa rescisória de 10% sobre as parcelas vincendas restantes.")
    
    doc.add_heading("3. Trancamento de Matrícula", level=2)
    doc.add_paragraph("O estudante pode solicitar o trancamento do curso por até 180 dias (6 meses) por motivos profissionais ou de saúde, pausando a contagem do tempo de acesso contratado.")
    doc.save(os.path.join(dir_path, "politica_reembolso_matriculas.docx"))

    # 5.3 FAQ Cursos e Certificados (.json)
    faq_data = {
        "dominio": "Plataforma Educativa / Escola Online",
        "instituicao": "Alura Tech Academy",
        "faq_cursos_certificados": [
            {
                "id": "FAQ-EDU-01",
                "pergunta": "Os certificados são reconhecidos e válidos para horas complementares?",
                "resposta": "Sim! Nossos certificados são emitidos em conformidade com a legislação de Cursos Livres (Lei nº 9.394/96) e possuem carga horária descrita e autenticação por QR Code, aceitos como atividades complementares universitárias."
            },
            {
                "id": "FAQ-EDU-02",
                "pergunta": "Por quanto tempo tenho acesso aos cursos?",
                "resposta": "O plano anual garante 12 meses de acesso ilimitado com todas as atualizações de conteúdo incluídas. Cursos individuais avulsos possuem acesso vitalício."
            },
            {
                "id": "FAQ-EDU-03",
                "pergunta": "Como tirar dúvidas com os professores?",
                "resposta": "Você pode postar sua dúvida diretamente no fórum abaixo de cada aula. Nossa equipe pedagógica e instrutores respondem em um tempo médio inferior a 4 horas úteis."
            }
        ]
    }
    with open(os.path.join(dir_path, "faq_cursos_certificados.json"), "w", encoding="utf-8") as f:
        json.dump(faq_data, f, ensure_ascii=False, indent=2)

    # 5.4 Guia de uso da plataforma (.md)
    md_content = """# Guia de Uso da Plataforma — Alura Tech Academy

## 1. Requisitos Técnicos
- **Navegadores Homologados**: Google Chrome (versão 100+), Mozilla Firefox, Microsoft Edge e Safari.
- **Conexão de Internet**: Velocidade mínima recomendada de 5 Mbps para reprodução fluida em Full HD (1080p).

## 2. Recursos do Player de Aulas
- Controle de velocidade de reprodução (0.5x, 1.0x, 1.25x, 1.5x, 1.75x e 2.0x).
- Legendas em português e transcrição textual completa de cada aula.
- Atalhos de teclado: Barra de espaço (Play/Pause), Teclas de seta (Avançar/Retroceder 5s).
- Download de arquivos de código-fonte e slides complementares na aba 'Materiais de Apoio'.

## 3. Comunidade e Networking
Alunos matriculados ganham acesso automático ao servidor exclusivo no Discord para networking, grupos de estudo e canais de vagas de emprego.
"""
    with open(os.path.join(dir_path, "guia_uso_plataforma.md"), "w", encoding="utf-8") as f:
        f.write(md_content)

    # 5.5 Programa de bolsas e afiliados (.csv)
    with open(os.path.join(dir_path, "programa_bolsas_afiliados.csv"), "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["modalidade", "tipo", "percentual_desconto_ou_comissao", "requisitos_principais", "vagas_por_ciclo"])
        dados = [
            ["Bolsa Mérito Acadêmico", "Bolsa de Estudos", "50%", "Nota acima de 85 no teste de lógica da plataforma", 100],
            ["Bolsa Diversidade & Inclusão", "Bolsa de Estudos", "100%", "Renda familiar per capita até 1.5 salário mínimo e autodeclaração", 50],
            ["Bolsa Convênio Corporativo", "Desconto Parceria", "30%", "Vínculo empregatício comprovado com empresas parceiras", "Ilimitado"],
            ["Programa de Afiliados Alura Tech", "Comissão Parceiro", "20%", "Indicação de novos alunos via link exclusivo com pagamento mensal", "Ilimitado"],
        ]
        writer.writerows(dados)


# ==========================================
# 6. FINTECH / BANCO DIGITAL
# ==========================================
def gerar_fintech(dir_path: str):
    os.makedirs(dir_path, exist_ok=True)

    # 6.1 Privacidade e protecao de dados (.pdf)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", "B", 15)
    pdf.cell(0, 10, "Politica de Privacidade e Sigilo Bancario - NovaBank Fintech", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(3)
    pdf.set_font("helvetica", "", 10)
    regras_fin = [
        "1. Conformidade Regulatoria e Sigilo Bancario:",
        "   - Em conformidade com a Lei Complementar no 105/2001 (Sigilo Bancario) e com a LGPD.",
        "   - Infraestrutura em conformidade com as resolucoes do Banco Central do Brasil (Bacen).",
        "",
        "2. Criptografia e Armazenamento:",
        "   - Dados de cartao e transacoes protegidos por criptografia de ponta a ponta AES-256 e certificacao PCI-DSS.",
        "   - Autenticacao biometrica facial (FaceMatch) obrigatoria para transacoes de alto valor.",
        "",
        "3. Compartilhamento via Open Finance:",
        "   - O compartilhamento de dados bancarios com outras instituicoes financeiras ocorre SOMENTE mediante",
        "     consentimento formal e explicito do cliente, com validade maxima de 12 meses, revogavel a qualquer momento.",
    ]
    for r in regras_fin:
        pdf.multi_cell(0, 5, r, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.output(os.path.join(dir_path, "privacidade_protecao_dados.pdf"))

    # 6.2 Termos e condicoes de uso (.md)
    md_content = """# Termos e Condições de Uso da Conta Digital — NovaBank

Última revisão: Fevereiro de 2026.

## 1. Condições para Abertura de Conta
- A conta digital NovaBank é destinada a pessoas físicas maiores de 18 anos (ou maiores de 16 anos emancipados) residentes no Brasil.
- É obrigatório possuir CPF regular perante a Receita Federal e enviar documento oficial com foto e selfie biométrica.

## 2. Rendimento Automático da Conta
- Todo o saldo mantido na conta digital NovaBank rende automaticamente **100% do CDI** em todos os dias úteis.
- O saldo possui liquidez imediata diária para transferências, pagamentos e saques a qualquer momento.

## 3. Cartões de Débito e Crédito
- Abertura de conta inclui cartão virtual gratuito instantâneo para compras online.
- O cartão físico múltiplo (débito e crédito) com bandeira Mastercard Internacional não possui taxa de anuidade na categoria Gold.
"""
    with open(os.path.join(dir_path, "termos_condicoes_uso.md"), "w", encoding="utf-8") as f:
        f.write(md_content)

    # 6.3 FAQ Transacoes e Limites (.json)
    faq_data = {
        "dominio": "Fintech / Banco Digital",
        "instituicao": "NovaBank Pagamentos S.A.",
        "faq_transacoes_limites": [
            {
                "id": "FAQ-FIN-01",
                "pergunta": "Quais são os limites de transferência Pix?",
                "resposta": "O limite padrão de Pix diurno (das 06h00 às 20h00) é de R$ 5.000,00 por dia. O limite noturno (das 20h00 às 06h00) é de R$ 1.000,00 por determinação de segurança do Banco Central. Solicitações de aumento de limite são avaliadas em até 24 a 48 horas úteis."
            },
            {
                "id": "FAQ-FIN-02",
                "pergunta": "O Pix ou TED possuem cobrança de tarifas?",
                "resposta": "Não! Todas as transferências Pix e TED são 100% gratuitas e ilimitadas para pessoas físicas."
            },
            {
                "id": "FAQ-FIN-03",
                "pergunta": "Como funcionam os saques em dinheiro?",
                "resposta": "Os saques podem ser realizados em qualquer caixa eletrônico da rede Banco24Horas. Todo cliente possui 4 saques gratuitos por mês. A partir do 5º saque no mesmo mês, é cobrada a tarifa de R$ 6,50 por operação."
            }
        ]
    }
    with open(os.path.join(dir_path, "faq_transacoes_limites.json"), "w", encoding="utf-8") as f:
        json.dump(faq_data, f, ensure_ascii=False, indent=2)

    # 6.4 Seguranca e prevencao de fraudes (.docx)
    doc = DocxDocument()
    doc.add_heading("Política de Segurança da Informação e Prevenção de Fraudes — NovaBank", level=1)
    
    doc.add_heading("1. Mecanismos de Autenticação e Antifraude", level=2)
    doc.add_paragraph("O NovaBank utiliza inteligência artificial para análise de risco em tempo real em 100% das transações. Qualquer comportamento atípico aciona o bloqueio preventivo da transação e notificação imediata via notificação push no app.")
    
    doc.add_heading("2. Mecanismo Especial de Devolução (MED) do Pix", level=2)
    doc.add_paragraph("Em caso de golpe ou suspeita de fraude via Pix, o cliente deve registrar a contestação no app em até 80 dias da transação. Os recursos são bloqueados cautelarmente na conta recebedora para análise do caso e estorno.")
    
    doc.add_heading("3. Central de Emergência 24 Horas", level=2)
    doc.add_paragraph("Linha telefônica dedicada 24 horas por dia, 7 dias por semana para bloqueio imediato de cartões e conta em caso de perda, roubo ou invasão de celular: 0800-999-NOVA (0800-999-6682).")
    doc.save(os.path.join(dir_path, "seguranca_prevencao_fraudes.docx"))

    # 6.5 Tarifas e comissoes (.xlsx)
    df_tarifas = pd.DataFrame([
        {"Servico": "Abertura e Manutenção de Conta Digital", "Tarifa_Pessoa_Fisica": "R$ 0,00 (Gratuito)", "Periodicidade": "Mensal"},
        {"Servico": "Transferência Pix (Envio e Recebimento)", "Tarifa_Pessoa_Fisica": "R$ 0,00 (Gratuito)", "Periodicidade": "Por transação"},
        {"Servico": "Transferência TED para outros bancos", "Tarifa_Pessoa_Fisica": "R$ 0,00 (Gratuito)", "Periodicidade": "Por transação"},
        {"Servico": "Emissão de Cartão Virtual", "Tarifa_Pessoa_Fisica": "R$ 0,00 (Gratuito)", "Periodicidade": "Ilimitado"},
        {"Servico": "Saque na Rede Banco24Horas (Até 4 saques/mês)", "Tarifa_Pessoa_Fisica": "R$ 0,00 (Gratuito)", "Periodicidade": "Mensal"},
        {"Servico": "Saque Adicional na Rede Banco24Horas (A partir do 5º)", "Tarifa_Pessoa_Fisica": "R$ 6,50", "Periodicidade": "Por saque adicional"},
        {"Servico": "Emissão de 2ª via de Cartão Físico (Perda/Dano)", "Tarifa_Pessoa_Fisica": "R$ 15,00", "Periodicidade": "Por emissão"},
        {"Servico": "Anuidade Cartão Mastercard Gold", "Tarifa_Pessoa_Fisica": "R$ 0,00 (Isento)", "Periodicidade": "Anual"},
    ])
    with pd.ExcelWriter(os.path.join(dir_path, "tarifas_e_comissoes.xlsx"), engine="openpyxl") as writer:
        df_tarifas.to_excel(writer, sheet_name="Tarifas_NovaBank", index=False)

    # 6.6 Apresentacao institucional em PowerPoint (.pptx)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide_layout = prs.slide_layouts[1]
    
    # Slide 1
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "NovaBank — Governança e Segurança Digital"
    tf1 = slide1.shapes.placeholders[1].text_frame
    tf1.text = "Pilares de Atuação da Fintech:"
    p = tf1.add_paragraph()
    p.text = "• 100% Digital e Seguro: Regulado pelo Banco Central do Brasil."
    p = tf1.add_paragraph()
    p.text = "• Rendimento Automático: 100% do CDI com liquidez diária."
    p = tf1.add_paragraph()
    p.text = "• Zero Anuidade: Cartões internacionais sem tarifas abusivas."

    # Slide 2
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "NovaBank — Atendimento e Proteção ao Cliente"
    tf2 = slide2.shapes.placeholders[1].text_frame
    tf2.text = "Compromissos Operacionais:"
    p = tf2.add_paragraph()
    p.text = "• Prevenção de Fraudes com IA em tempo real e MED Pix em até 80 dias."
    p = tf2.add_paragraph()
    p.text = "• Central 24h pelo 0800-999-NOVA para emergências."
    p = tf2.add_paragraph()
    p.text = "• 4 saques gratuitos por mês na rede Banco24Horas."

    prs.save(os.path.join(dir_path, "apresentacao_institucional.pptx"))


def main():
    print(f"Gerando documentação do Challenge em: {DOCS_DIR}")
    gerar_ecommerce(os.path.join(DOCS_DIR, "ecommerce"))
    gerar_saas(os.path.join(DOCS_DIR, "saas"))
    gerar_logistica(os.path.join(DOCS_DIR, "logistica"))
    gerar_saude(os.path.join(DOCS_DIR, "saude"))
    gerar_educacao(os.path.join(DOCS_DIR, "educacao"))
    gerar_fintech(os.path.join(DOCS_DIR, "fintech"))
    print("Todos os documentos dos 6 segmentos foram gerados com sucesso!")


if __name__ == "__main__":
    main()