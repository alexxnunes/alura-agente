import os

import pandas as pd
from docx import Document
from pptx import Presentation

from src.loaders import SUPPORTED_EXTS, load_document

DOCS = os.path.join(os.path.dirname(__file__), "..", "data", "docs")


def test_suporta_os_oito_formatos_do_desafio():
    esperados = {".pdf", ".docx", ".xlsx", ".pptx", ".md", ".csv", ".json", ".html"}
    assert esperados.issubset(set(SUPPORTED_EXTS))


def test_carrega_markdown():
    texto = load_document(os.path.join(DOCS, "stack_tecnologica.md"))
    assert "FastAPI" in texto and "PostgreSQL" in texto


def test_carrega_csv():
    texto = load_document(os.path.join(DOCS, "vendas_produtos_2023_2025.csv"))
    assert "Smartphone Zenith Pro" in texto and "dezembro" in texto


def test_carrega_json():
    texto = load_document(os.path.join(DOCS, "produtos_catalogo.json"))
    assert "Notebook Ultra X1" in texto and "230" in texto


def test_carrega_html():
    texto = load_document(os.path.join(DOCS, "comunicado_interno.html"))
    assert "auxilio academia" in texto


def test_carrega_pdf():
    texto = load_document(os.path.join(DOCS, "politicas_rh.pdf"))
    assert "Home office" in texto


def test_carrega_word_excel_e_powerpoint(tmp_path):
    caminho_docx = tmp_path / "manual.docx"
    documento = Document()
    documento.add_paragraph("Manual de onboarding")
    documento.save(caminho_docx)

    caminho_xlsx = tmp_path / "indicadores.xlsx"
    pd.DataFrame([{"indicador": "NPS", "meta": 75}]).to_excel(
        caminho_xlsx, sheet_name="Metas", index=False
    )

    caminho_pptx = tmp_path / "roadmap.pptx"
    apresentacao = Presentation()
    slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[5])
    slide.shapes.title.text = "Roadmap estratégico 2026"
    apresentacao.save(caminho_pptx)

    assert "Manual de onboarding" in load_document(str(caminho_docx))
    assert "NPS" in load_document(str(caminho_xlsx))
    assert "Roadmap estratégico 2026" in load_document(str(caminho_pptx))


def test_extensao_desconhecida_levanta_erro(tmp_path):
    path = tmp_path / "arquivo.txt"
    path.write_text("conteudo", encoding="utf-8")
    try:
        load_document(str(path))
    except ValueError as exc:
        assert ".txt" in str(exc)
    else:
        raise AssertionError("deveria levantar ValueError")


def test_carrega_documentos_dos_segmentos():
    # Fintech
    texto_fin_pdf = load_document(os.path.join(DOCS, "fintech", "privacidade_protecao_dados.pdf"))
    assert "Sigilo Bancario" in texto_fin_pdf or "LGPD" in texto_fin_pdf
    texto_fin_xlsx = load_document(os.path.join(DOCS, "fintech", "tarifas_e_comissoes.xlsx"))
    assert "Banco24Horas" in texto_fin_xlsx

    # SaaS
    texto_saas_md = load_document(os.path.join(DOCS, "saas", "base_conhecimento.md"))
    assert "CloudSync" in texto_saas_md
    texto_saas_docx = load_document(os.path.join(DOCS, "saas", "politica_privacidade.docx"))
    assert "TLS 1.3" in texto_saas_docx

    # E-commerce
    texto_ec_docx = load_document(os.path.join(DOCS, "ecommerce", "guia_envios_entregas.docx"))
    assert "Sedex" in texto_ec_docx

    # Logistica
    texto_log_md = load_document(os.path.join(DOCS, "logistica", "rastreamento_pedidos.md"))
    assert "TransLogística" in texto_log_md

    # Saúde
    texto_saude_xlsx = load_document(os.path.join(DOCS, "saude", "convenios_coberturas.xlsx"))
    assert "Unimed" in texto_saude_xlsx

    # Educação
    texto_edu_csv = load_document(os.path.join(DOCS, "educacao", "programa_bolsas_afiliados.csv"))
    assert "Bolsa Mérito" in texto_edu_csv

